from pptx import Presentation

def write_pptx(data, filename="mini_deck.pptx"):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title & Content

    # ✅ Slide 1: Company Overview
    slide1 = prs.slides.add_slide(slide_layout)
    slide1.shapes.title.text = "Company Overview"
    co = data["company_overview"]
    overview_text = (
        f"Name: {co['name']}\n"
        f"Sector: {co['sector']}\n"
        f"HQ: {co['headquarters']}\n\n"
        f"{co['business_model']}"
    )
    slide1.shapes.placeholders[1].text = overview_text

    # ✅ Slide 2: Key Financials (updated keys)
    slide2 = prs.slides.add_slide(slide_layout)
    slide2.shapes.title.text = "Key Financials"
    kf = data["key_financials"]
    kf_text = (
        f"- Revenue (Actual): {kf['revenue_actual']}\n"
        f"- Revenue (Forecast): {kf['revenue_forecast']}\n"
        f"- EBITDA (Actual): {kf['ebitda_actual']}\n"
        f"- EBITDA (Forecast): {kf['ebitda_forecast']}\n"
        f"- Recent CAGR: {kf['recent_cagr']}"
    )
    slide2.shapes.placeholders[1].text = kf_text

    # ✅ Slide 3: Investment Rationale
    slide3 = prs.slides.add_slide(slide_layout)
    slide3.shapes.title.text = "Investment Rationale"
    rationale = "\n".join(f"- {pt}" for pt in data["investment_highlights"][:5])
    slide3.shapes.placeholders[1].text = rationale

    prs.save(filename)
    print(f"✅ PPTX saved: {filename}")
